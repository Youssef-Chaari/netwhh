import os
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Erreur : La bibliothèque 'python-pptx' n'est pas installée.")
    print("Veuillez lancer : pip install python-pptx")
    exit()

def create_presentation():
    prs = Presentation()

    # Couleurs et styles
    # (Par défaut, on utilise un template simple et propre)

    slides_content = [
        ["Présentation Finale : Assurance Qualité BiProject", "Binôme : Youssef Chaari & Mohamed Aziz Zouari\nProjet de Fin d'Étude — QA & Automatisation"],
        ["Objectifs de la Mission", "- Garantir la stabilité du flux d'achat (E2E)\n- Identifier les failles critiques (Sécurité)\n- Automatiser la régression (Pyramide des tests)\n- Assurer la conformité académique"],
        ["Stratégie et Pyramide des Tests", "- Tests Unitaires : Logique d'auth (.NET)\n- Tests d'Intégration : Validation API & Contrôle d'accès\n- Tests Système : Parcours clients complets (Selenium)\n- Analyse Statique : Sécurité & Structure"],
        ["Analyse Statique & Sécurité", "- 3 Anomalies majeures détectées (IDOR, Trust Boundary, Logging)\n- Traçabilité : Matrice reliant 100% des exigences aux cas de test"],
        ["Tests d'Intégration & Preuve de Bug", "- Le cas CT-20 : Validation dynamique de la faille IDOR\n- Résultat : Échec révélateur (Accès 200 OK non autorisé)\n- Correction & Régression : Test CT-18 validé"],
        ["Automatisation Selenium (Principal)", "- Architecture POM (Page Object Model)\n- Robustesse : BasePage & Explicit Waits\n- Observabilité : Captures d'écran automatiques sur échec"],
        ["Bonus : Robot Framework (Prototype)", "- Approche Keyword-Driven : Lisibilité maximale\n- Non exécuté officiellement (Preuve de concept)\n- Objectif : Polyvalence technologique (Pytest vs Robot)"],
        ["Indicateurs & Bilan Final", "- 20 Cas de tests couverts\n- 85% de taux de réussite (17/20 Pass)\n- Livrables : Catalogue, Rapports, Screenshots"],
        ["Conclusion & Questions", "Projet Conforme, Audité et Maintenable.\nMerci de votre attention !"]
    ]

    for title_text, body_text in slides_content:
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = title_text
        content.text = body_text

    output_path = "Presentation_QA_BiProject.pptx"
    prs.save(output_path)
    print(f"Félicitations ! Votre présentation a été créée : {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
